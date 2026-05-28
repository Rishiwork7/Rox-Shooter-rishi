import tkinter as tk
import customtkinter as ctk
import multiprocessing
import threading
import queue
import asyncio
import os
import shutil
import string
import re
import random
import tempfile
from datetime import datetime, timedelta
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PILImage
import xlsxwriter
import subprocess
import sys
import uuid
import pypdfium2 as pdfium
import requests
from supabase import create_client, Client
import base64
from io import BytesIO
def resource_path(relative_path):
    """ Get absolute path to resource """
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base_path, relative_path)

# --- SUPABASE CREDENTIALS ---
SUPABASE_URL = "https://syzmaecfeiltzrtmlgoq.supabase.co"
SUPABASE_KEY = "sb_publishable_IPDIsxft6C9RRy4s9EPgOQ_rXVaC8N-"

ctk.set_appearance_mode("light")
ctk.set_default_color_theme("blue")

# ── Light Theme Color Palette ──
COLOR = {
    "bg":           "#F5F7FA",
    "surface":      "#FFFFFF",
    "surface_alt":  "#EEF1F6",
    "border":       "#DDE2EA",
    "text":         "#1E293B",
    "text_sec":     "#64748B",
    "primary":      "#4F6DF5",
    "primary_hov":  "#3B5BDB",
    "success":      "#22C55E",
    "success_hov":  "#16A34A",
    "danger":       "#EF4444",
    "danger_hov":   "#DC2626",
    "warn":         "#F59E0B",
    "warn_hov":     "#D97706",
    "muted":        "#94A3B8",
    "muted_hov":    "#64748B",
    "accent":       "#8B5CF6",
    "accent_hov":   "#7C3AED",
    "info":         "#0EA5E9",
    "info_hov":     "#0284C7",
    "input_bg":     "#F1F5F9",
    "input_border": "#CBD5E1",
    "log_bg":       "#F8FAFC",
    "log_text":     "#334155",
    "tag_card":     "#F1F5F9",
}

class TextParser:
    def __init__(self):
        # Realistic placeholder names for USA tags
        self.first_names = ["James", "Robert", "John", "Michael", "David", "William", "Richard", "Joseph", "Thomas", "Charles", "Christopher", "Daniel", "Matthew", "Anthony", "Mark", "Donald", "Steven", "Paul", "Andrew", "Joshua", "Mary", "Patricia", "Jennifer", "Linda", "Elizabeth", "Barbara", "Susan", "Jessica", "Sarah", "Karen", "Lisa", "Nancy", "Betty", "Margaret", "Sandra", "Ashley", "Kimberly", "Emily", "Donna", "Michelle"]
        self.last_names = ["Smith", "Johnson", "Williams", "Brown", "Jones", "Garcia", "Miller", "Davis", "Rodriguez", "Martinez", "Hernandez", "Lopez", "Gonzalez", "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "Martin", "Lee", "Perez", "Thompson", "White", "Harris", "Sanchez", "Clark", "Ramirez", "Lewis", "Robinson", "Walker", "Young", "Allen", "King", "Wright", "Scott", "Torres", "Nguyen", "Hill", "Flores"]

    def generate_random_string(self, length):
        return ''.join(random.choices(string.ascii_uppercase + string.digits, k=length))

    def generate_word(self, length):
        return ''.join(random.choices(string.ascii_uppercase, k=length))

    def generate_invoice(self, length=12):
        # Format like 22250174NCN8
        num_len = max(1, length - 4)
        num_part = ''.join(random.choices(string.digits, k=num_len))
        alpha_part = ''.join(random.choices(string.ascii_uppercase, k=4))
        return num_part + alpha_part

    def generate_number(self, length):
        return ''.join(random.choices(string.digits, k=length))

    def generate_usa_name(self):
        return f"{random.choice(self.first_names)} {random.choice(self.last_names)}"

    def generate_usa_rnd(self):
        fn = random.choice(self.first_names)
        ln = random.choice(self.last_names)
        formats = [
            f"{fn[0]}.{random.choice(string.ascii_uppercase)}. {ln}",
            f"{ln}, {fn} {random.choice(string.ascii_uppercase)}.",
            f"{fn} {ln}",
            f"{ln}, {fn}"
        ]
        return random.choice(formats)

    def parse(self, text, recipient_email, tfn_number):
        if not text:
            return ""

        username = recipient_email.split('@')[0] if '@' in recipient_email else recipient_email

        def replacer(match):
            tag = match.group(1).lower()
            val = match.group(2)
            length = int(val) if val else None

            if tag == "random":
                return self.generate_random_string(length or 6)
            elif tag == "word":
                return self.generate_word(length or 6)
            elif tag == "invoice_no":
                return self.generate_invoice(length or 12)
            elif tag in ["rand", "number"]:
                return self.generate_number(length or 6)
            elif tag in ["mail", "email"]:
                return recipient_email
            elif tag == "n_mail":
                return username
            elif tag == "date":
                return datetime.now().strftime("%d-%m-%y")
            elif tag == "tfn":
                return tfn_number
            elif tag == "number6":
                return self.generate_number(6)
            elif tag == "random8":
                return self.generate_random_string(8)
            elif tag == "word3":
                return self.generate_word(3)
            elif tag == "word4":
                return self.generate_word(4)
            elif tag == "digi_no":
                return str(uuid.uuid4())
            elif tag == "usa_name":
                return self.generate_usa_name()
            elif tag == "usa_rnd":
                return self.generate_usa_rnd()
            elif tag == "str_ing":
                return self.generate_word(random.choice([5, 11]))
            elif tag == "mail_rnd":
                return f"{username}_{self.generate_number(8)}"
            elif tag == "fn_usa":
                return random.choice(self.first_names)
            elif tag == "charge":
                return str(random.randint(10, 99))
            
            return match.group(0)

        # Match tags including those with underscores and numbers
        supported_tags = [
            "random", "word", "invoice_no", "rand", "number", "mail", "email", "date", "tfn",
            "n_mail", "number6", "random8", "digi_no", "usa_name", "usa_rnd", "str_ing", 
            "mail_rnd", "fn_usa", "charge", "word3", "word4"
        ]
        pattern = r"\$(" + "|".join(supported_tags) + r")(?:[\(\[])?(\d+)?(?:[\)\]])?"
        text = re.sub(pattern, replacer, text, flags=re.IGNORECASE)

        return text

class Converter:
    def __init__(self, log_callback):
        self.log = log_callback
        self.delays = {}
        # Use system temp directory instead of relative path
        # This ensures compatibility with PyInstaller .exe on Windows
        self.temp_dir = os.path.join(tempfile.gettempdir(), "rox_shooter_attachments")
        if os.path.exists(self.temp_dir):
            try:
                shutil.rmtree(self.temp_dir)
            except Exception as e:
                self.log(f"Warning: Could not clean temp directory {self.temp_dir}: {e}")
        try:
            os.makedirs(self.temp_dir, exist_ok=True)
        except Exception as e:
            self.log(f"Error creating temp directory: {e}")
            
        self._playwright = None
        self._browser = None
        self._browser_lock = None

    async def _get_browser(self):
        if self._browser_lock is None:
            self._browser_lock = asyncio.Lock()
        async with self._browser_lock:
            if self._browser is None:
                self._playwright = await async_playwright().start()
                self._browser = await self._playwright.chromium.launch(headless=True, channel="chrome")
            return self._browser

    async def html_to_pdf(self, html_content, filename="attachment.pdf"):
        """Convert HTML to PDF using persistent Chrome."""
        path = os.path.join(self.temp_dir, filename)
        for attempt in range(1, 4):
            page = None
            try:
                browser = await self._get_browser()
                page = await browser.new_page()
                await page.set_content(html_content, wait_until="networkidle")
                await page.pdf(path=path, format="A4")
                await page.close()
                return path
            except Exception as e:
                self.log(f"PDF Conversion attempt {attempt}/3 failed: {e}")
                if page:
                    try: await page.close()
                    except: pass
                if attempt < 3:
                    await asyncio.sleep(self.delays.get('retry', 1.5))
        return None

    async def html_to_image(self, html_content, filename="attachment.jpg"):
        """Convert HTML to Image using persistent Chrome."""
        path = os.path.join(self.temp_dir, filename)
        for attempt in range(1, 4):
            page = None
            try:
                browser = await self._get_browser()
                page = await browser.new_page()
                await page.set_content(html_content, wait_until="networkidle")
                await page.screenshot(path=path, type="jpeg", quality=60, full_page=True)
                await page.close()
                return path
            except Exception as e:
                self.log(f"Image Conversion attempt {attempt}/3 failed: {e}")
                if page:
                    try: await page.close()
                    except: pass
                if attempt < 3:
                    await asyncio.sleep(self.delays.get('retry', 1.5))
        return None


    async def html_to_image_pptx(self, html_content, filename="attachment_img.pptx"):
        """HTML -> Image -> PPTX Slide."""
        path = os.path.join(self.temp_dir, filename)
        unique_id = uuid.uuid4().hex[:8]
        img_name = f"temp_pptx_{unique_id}.jpg"
        img_path = os.path.join(self.temp_dir, img_name)
        try:
            await self.html_to_image(html_content, img_name)
            prs = Presentation()
            # Set to A4 dimensions (approx)
            prs.slide_width = Inches(8.5)
            prs.slide_height = Inches(11)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Fill the slide
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            prs.save(path)
            if os.path.exists(img_path): os.remove(img_path)
            return path
        except Exception as e:
            self.log(f"Image then PPTX Error: {e}")
            return None

    async def html_to_image_pdf(self, html_content, filename="attachment_img.pdf"):
        """HTML -> Image -> PDF."""
        path = os.path.join(self.temp_dir, filename)
        unique_id = uuid.uuid4().hex[:8]
        img_name = f"temp_pdf_{unique_id}.jpg"
        img_path = os.path.join(self.temp_dir, img_name)
        try:
            await self.html_to_image(html_content, img_name)
            img = PILImage.open(img_path)
            pdf_img = img.convert('RGB')
            pdf_img.save(path)
            if os.path.exists(img_path): os.remove(img_path)
            return path
        except Exception as e:
            self.log(f"Image then PDF Error: {e}")
            return None

    async def html_to_image_xls(self, html_content, filename="attachment_img.xlsx"):
        """HTML -> Image -> XLS."""
        path = os.path.join(self.temp_dir, filename)
        unique_id = uuid.uuid4().hex[:8]
        img_name = f"temp_xls_{unique_id}.jpg"
        img_path = os.path.join(self.temp_dir, img_name)
        try:
            await self.html_to_image(html_content, img_name)
            workbook = xlsxwriter.Workbook(path)
            worksheet = workbook.add_worksheet()
            worksheet.insert_image('B2', img_path, {'x_scale': 0.5, 'y_scale': 0.5})
            workbook.close()
            if os.path.exists(img_path): os.remove(img_path)
            return path
        except Exception as e:
            self.log(f"Image then XLS Error: {e}")
            return None

    async def html_to_pdf_pptx(self, html_content, filename="attachment_pdf.pptx"):
        """HTML -> PDF -> PPTX (Multi-page)."""
        path = os.path.join(self.temp_dir, filename)
        unique_id = uuid.uuid4().hex[:8]
        pdf_name = f"temp_pdf_{unique_id}.pdf"
        pdf_path = os.path.join(self.temp_dir, pdf_name)
        try:
            # 1. Generate PDF
            res = await self.html_to_pdf(html_content, pdf_name)
            if not res: return None
            
            # 2. Convert PDF to Images
            pdf = pdfium.PdfDocument(pdf_path)
            prs = Presentation()
            # Set to A4 dimensions
            prs.slide_width = Inches(8.5)
            prs.slide_height = Inches(11)
            temp_images = []
            
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                bitmap = page.render(scale=2) # Higher DPI
                pil_image = bitmap.to_pil()
                
                img_path = os.path.join(self.temp_dir, f"page_{page_num}_{unique_id}.png")
                pil_image.save(img_path)
                temp_images.append(img_path)
                
                # 3. Add to PPTX
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # Stretch to fill the A4 slide
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
            prs.save(path)
            pdf.close()
            
            # Cleanup
            for img in temp_images:
                if os.path.exists(img): os.remove(img)
            if os.path.exists(pdf_path): os.remove(pdf_path)
            return path
        except Exception as e:
            self.log(f"PDF then PPTX Error: {e}")
            return None

    async def html_to_pdf_xls(self, html_content, filename="attachment_pdf.xlsx"):
        """HTML -> PDF -> Excel (Multi-page)."""
        path = os.path.join(self.temp_dir, filename)
        unique_id = uuid.uuid4().hex[:8]
        pdf_name = f"temp_pdf_{unique_id}.pdf"
        pdf_path = os.path.join(self.temp_dir, pdf_name)
        try:
            # 1. Generate PDF
            res = await self.html_to_pdf(html_content, pdf_name)
            if not res: return None
            
            # 2. Convert PDF to Images
            pdf = pdfium.PdfDocument(pdf_path)
            workbook = xlsxwriter.Workbook(path)
            temp_images = []
            worksheet = workbook.add_worksheet("Document")
            current_row = 0
            
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                # स्केल 1 (नॉर्मल DPI) करने से पिक्सल 4 गुना कम हो जाएंगे
                bitmap = page.render(scale=1)
                pil_image = bitmap.to_pil()
                
                # PNG की जगह JPEG (Quality=60) में सेव करेंगे
                img_path = os.path.join(self.temp_dir, f"page_{page_num}_{unique_id}.jpg")
                pil_image.convert('RGB').save(img_path, format="JPEG", quality=60, optimize=True)
                temp_images.append(img_path)
                
                # 3. Add to XLS (सभी पेजों को एक ही शीट में नीचे-नीचे लगाएंगे)
                worksheet.insert_image(current_row, 0, img_path, {'x_scale': 1.0, 'y_scale': 1.0})
                
                # अगली इमेज को नीचे खिसकाने के लिए row कैलकुलेट करेंगे (1 row = ~20 pixels)
                current_row += int(pil_image.height / 19) + 1
            
            workbook.close()
            pdf.close()
            
            # Cleanup
            for img in temp_images:
                if os.path.exists(img): os.remove(img)
            if os.path.exists(pdf_path): os.remove(pdf_path)
            return path
        except Exception as e:
            self.log(f"PDF then XLS Error: {e}")
            return None

def get_public_ip():
    """Fetches the current public IP address using reliable APIs."""
    try:
        response = requests.get('https://api.ipify.org', timeout=5)
        return response.text.strip()
    except Exception:
        try:
            response = requests.get('https://ifconfig.me/ip', timeout=5)
            return response.text.strip()
        except Exception:
            return None

def obfuscate_and_attach(msg_object, file_path, file_name):
    """
    Polymorphic Padding and MIME Hardening Engine.
    Mutates XLSX hash dynamically and constructs hardened MIMEBase attachment.
    """
    try:
        from email.mime.base import MIMEBase
        from email import encoders
        import os
        import random
        import string
        import uuid
        
        if not os.path.exists(file_path):
            return None

        # 2. POLYMORPHIC PADDING
        # Read the .xlsx file in binary mode
        with open(file_path, 'rb') as f:
            data = f.read()

        # Append 15 to 60 bytes of random characters to the end of binary data
        padding_size = random.randint(15, 60)
        padding_chars = ''.join(random.choices(string.ascii_letters + string.digits, k=padding_size))
        mutated_data = data + padding_chars.encode('utf-8')

        # Write mutated data back to path so that Playwright can also upload it with updated hash
        with open(file_path, 'wb') as f:
            f.write(mutated_data)

        # 3. MIME HARDENING
        # Create MIMEBase object
        mime_part = MIMEBase('application', 'octet-stream')
        mime_part.set_payload(mutated_data)
        encoders.encode_base64(mime_part)
        
        # Override default MIME boundary formats
        random_str = ''.join(random.choices(string.ascii_lowercase + string.digits, k=8))
        custom_filename = f"{random_str}_invoice.xlsx"
        
        # Remove standard Python header fingerprints
        mime_part.replace_header('Content-Type', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        mime_part.add_header(
            'Content-Disposition',
            f'attachment; filename="{custom_filename}"'
        )
        
        # Override default MIME boundaries if parent msg_object exists to mimic standard Apple Mail / Outlook format
        if msg_object:
            outlook_boundary = f"----=_NextPart_000_{random.randint(1000, 9999):04d}_{uuid.uuid4().hex[:8].upper()}.{uuid.uuid4().hex[:8].upper()}"
            msg_object.set_boundary(outlook_boundary)

        return mime_part
    except Exception as e:
        return None

def create_stealth_mime_payload(sender_email, recipient_email, subject, html_body, attachment_path):
    """
    Advanced SMTP & MIME Transport Obfuscation Engine.
    Mimics human-sent corporate email (Outlook/Apple Mail) to maximize inbox placement.

    Evasion Layers:
    1. MIME-Type fallback: application/octet-stream (bypasses deep packet inspection)
    2. 72-char Base64 chunking (spoofs Outlook/Apple Mail transport signature)
    3. Entropy anchor: hidden corporate boilerplate HTML (inflates text-trust score)
    """
    import base64
    import random
    import string
    import uuid
    import os
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText
    from email.mime.base import MIMEBase
    from email import encoders

    try:
        # --- RANDOM FILENAME & IDs ---
        invoice_id = ''.join(random.choices(string.digits, k=8))
        custom_filename = f"Invoice_{invoice_id}.xlsx"
        outlook_boundary = (
            f"----=_NextPart_000_{random.randint(1000,9999):04d}"
            f"_{uuid.uuid4().hex[:8].upper()}"
            f".{uuid.uuid4().hex[:8].upper()}"
        )

        # --- BUILD MULTIPART MESSAGE ---
        msg = MIMEMultipart('mixed', boundary=outlook_boundary)
        msg['From']    = sender_email
        msg['To']      = recipient_email
        msg['Subject'] = subject
        # Spoof desktop client headers
        msg['X-Mailer']   = f'Microsoft Outlook 16.0.{random.randint(14000, 17000)}'
        msg['MIME-Version'] = '1.0'

        # --- LAYER 3: ENTROPY ANCHOR ---
        # Inject massive trusted corporate boilerplate in near-invisible text
        boilerplate = (
            "<div style='font-size:9px;color:#d3d3d3;font-family:Arial,sans-serif;margin-top:40px;'>"
            "<p>CONFIDENTIALITY NOTICE: This email message, including any attachments, is for the sole use of the "
            "intended recipient(s) and may contain confidential and privileged information. Any unauthorized review, "
            "use, disclosure or distribution is prohibited. If you are not the intended recipient, please contact the "
            "sender by reply email and destroy all copies of the original message and any attachments.</p>"
            "<p>PRIVACY POLICY: We are committed to protecting your privacy and the confidentiality of your personal "
            "information. The personal information we collect is used only for the purposes for which it was collected "
            "and in accordance with applicable laws. For a complete description of our privacy practices, please visit "
            "our Privacy Policy page on our corporate website. You may opt-out of receiving commercial electronic "
            "messages at any time by following the unsubscribe instructions included in this email.</p>"
            "<p>LEGAL DISCLAIMER: The information contained in this communication is intended solely for the "
            "use of the individual or entity to whom it is addressed and others authorized to receive it. It may "
            "contain confidential or legally privileged information. If you are not the intended recipient you are "
            "hereby notified that any disclosure, copying, distribution or taking any action in reliance on the "
            "contents of this information is strictly prohibited and may be unlawful. &copy; "
            f"{random.randint(2019,2024)} All rights reserved.</p>"
            "</div>"
        )
        full_html = html_body + boilerplate
        msg.attach(MIMEText(full_html, 'html', 'utf-8'))

        # --- LAYER 1 + 2: MIME-TYPE FALLBACK & 72-CHAR BASE64 CHUNKING ---
        if attachment_path and os.path.exists(attachment_path):
            with open(attachment_path, 'rb') as f:
                raw_data = f.read()

            # Standard base64 encodes at 76 chars; we force 72 to spoof Outlook/Apple Mail
            b64_raw = base64.b64encode(raw_data).decode('ascii')
            chunked_b64 = '\n'.join(
                b64_raw[i:i+72] for i in range(0, len(b64_raw), 72)
            )

            # Force application/octet-stream to bypass DPI MIME inspection
            attachment = MIMEBase('application', 'octet-stream')
            attachment.set_payload(chunked_b64)
            attachment.add_header('Content-Transfer-Encoding', 'base64')
            attachment.add_header(
                'Content-Disposition',
                f'attachment; filename="{custom_filename}"'
            )
            # Intentionally override Content-Type to drop xlsx fingerprint
            attachment.replace_header(
                'Content-Type', f'application/octet-stream; name="{custom_filename}"'
            )
            msg.attach(attachment)

        return msg

    except Exception:
        return None

class LoginWindow(ctk.CTk):
    def __init__(self, on_success):
        super().__init__()
        self.on_success = on_success
        
        # --- Supabase Config (Using Hardcoded Constants) ---
        try:
            self.supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
        except Exception:
            self.supabase = None

        # --- Window Configuration ---
        self.title("Rox-Shooter v1.1 Authorization")
        self.geometry("400x500")
        self.configure(fg_color=COLOR["bg"])
        self.resizable(False, False)
        
        # --- Set Window Icon ---
        self.logo_path = resource_path("rox-logo.png")
        try:
            if os.path.exists(self.logo_path):
                img = tk.PhotoImage(file=self.logo_path)
                self.iconphoto(False, img)
        except:
            pass

        # Center Window
        screen_w = self.winfo_screenwidth()
        screen_h = self.winfo_screenheight()
        x = (screen_w // 2) - (400 // 2)
        y = (screen_h // 2) - (500 // 2)
        self.geometry(f"+{x}+{y}")

        self.create_widgets()
        
    def create_widgets(self):
        # Logo in place of Lock emoji
        try:
            if os.path.exists(self.logo_path):
                logo_img = ctk.CTkImage(light_image=PILImage.open(self.logo_path), dark_image=PILImage.open(self.logo_path), size=(64, 64))
                ctk.CTkLabel(self, text="", image=logo_img).pack(pady=(60, 10))
            else:
                ctk.CTkLabel(self, text="🔒", font=("Inter", 60)).pack(pady=(60, 10))
        except:
            ctk.CTkLabel(self, text="🔒", font=("Inter", 60)).pack(pady=(60, 10))
            
        ctk.CTkLabel(self, text="Rox-Shooter v1.1", font=("Inter", 24, "bold"), text_color=COLOR["text"]).pack(pady=(0, 30))
        
        ctk.CTkLabel(self, text="User Identifier", font=("Inter", 13, "bold"), text_color=COLOR["text_sec"]).pack(anchor="w", padx=40)
        self.user_id_entry = ctk.CTkEntry(self, width=320, height=40, corner_radius=8, fg_color=COLOR["input_bg"], border_color=COLOR["input_border"])
        self.user_id_entry.pack(pady=(5, 20))
        
        ctk.CTkLabel(self, text="Security Password", font=("Inter", 13, "bold"), text_color=COLOR["text_sec"]).pack(anchor="w", padx=40)
        self.password_entry = ctk.CTkEntry(self, width=320, height=40, corner_radius=8, fg_color=COLOR["input_bg"], border_color=COLOR["input_border"], show="*")
        self.password_entry.pack(pady=(5, 20))
        
        self.lbl_status = ctk.CTkLabel(self, text="", font=("Inter", 12), text_color=COLOR["danger"])
        self.lbl_status.pack(pady=10)
        
        self.btn_login = ctk.CTkButton(self, text="SECURE LOGIN", width=320, height=50, corner_radius=25,
                                      fg_color=COLOR["primary"], hover_color=COLOR["primary_hov"],
                                      font=("Inter", 14, "bold"), command=self.on_login_click)
        self.btn_login.pack(pady=10)

    def on_login_click(self):
        threading.Thread(target=self.login_process, daemon=True).start()

    def login_process(self):
        u_id = self.user_id_entry.get().strip()
        pwd = self.password_entry.get().strip()
        
        if not u_id or not pwd:
            self.update_status("Please fill all fields.", COLOR["danger"])
            return

        self.update_status("Verifying IP & Credentials...", COLOR["info"])
        current_ip = get_public_ip()
        if not current_ip:
            self.update_status("Network Error: Could not fetch IP.", COLOR["danger"])
            return

        if not self.supabase:
            self.update_status("Configuration Error: Check .env", COLOR["danger"])
            return
            
        try:
            response = self.supabase.table("users").select("*").eq("user_id", u_id).execute()
            data = response.data
            
            if not data:
                # Try case-insensitive fallback
                response = self.supabase.table("users").select("*").ilike("user_id", u_id).execute()
                data = response.data

            if not data:
                self.update_status("Invalid User ID.", COLOR["danger"])
                return
            
            # 3. Validations
            user = data[0]
            if user.get("password") != pwd:
                self.update_status("Incorrect Password.", COLOR["danger"])
                return
            if user.get("status") != "Active":
                self.update_status("Account is Inactive. Contact Admin.", COLOR["danger"])
                return
            
            # --- IP Validation & Multi-Binding ---
            allowed_ips_str = str(user.get("allowed_ip", "")).strip()
            ip_limit = int(user.get("ip_limit", 1))
            
            # Clean up the list of IPs
            if allowed_ips_str.upper() == "PENDING":
                ips = []
            elif allowed_ips_str in ["*", "Any", "any", ""]:
                ips = ["*"]
            else:
                ips = [ip.strip() for ip in allowed_ips_str.split(",") if ip.strip()]

            if "*" in ips:
                # Bypass check if "Any" or "*" is set
                pass
            elif current_ip in ips:
                # IP already authorized
                pass
            else:
                # New IP detected - check limit
                if len(ips) < ip_limit:
                    ips.append(current_ip)
                    new_allowed_ips = ",".join(ips)
                    try:
                        self.update_status(f"Binding new IP ({len(ips)}/{ip_limit})...", COLOR["info"])
                        self.supabase.table("users").update({"allowed_ip": new_allowed_ips}).eq("user_id", user.get("user_id")).execute()
                        self.update_status("IP authorized successfully!", COLOR["success"])
                    except Exception as e:
                        self.update_status(f"IP Binding Error: {str(e)}", COLOR["danger"])
                        return
                else:
                    self.update_status(f"Login Error: IP limit reached ({ip_limit}/{ip_limit})", COLOR["danger"])
                    return
            
            self.update_status("Authorization Success!", COLOR["success"])
            self.after(1000, self.finish_login)
        except Exception as e:
            self.update_status(f"Auth Error: {str(e)}", COLOR["danger"])

    def update_status(self, text, color):
        self.after(0, lambda: self.lbl_status.configure(text=text, text_color=color))

    def finish_login(self):
        u_id = self.user_id_entry.get().strip()
        self.destroy()
        self.on_success(u_id)

class App(ctk.CTk):
    def __init__(self, user_id="Admin"):
        super().__init__()
        self.user_id = user_id

        # --- Window Configuration ---
        self.geometry("1240x820")
        self.minsize(960, 640)
        self.configure(fg_color=COLOR["bg"])

        # --- Set Window Icon ---
        self.logo_path = resource_path("rox-logo.png")
        try:
            if os.path.exists(self.logo_path):
                img = tk.PhotoImage(file=self.logo_path)
                self.iconphoto(False, img)
        except:
            pass

        # UI Scaling/State
        self.after(0, lambda: self.state('zoomed'))
        
        # --- Internal State & Threading ---
        self.log_queue = queue.Queue()
        self.loop = None
        self.thread = None
        self.is_blasting = False
        self.is_closing = False
        self.contexts = {} # {id: {"context": context, "page": page, "frame": frame}}
        self.parser = TextParser()
        self.converter = Converter(self.log)
        
        # --- Background Loop ---
        self.loop = asyncio.new_event_loop()
        self.thread = threading.Thread(target=self.run_background_loop, daemon=True)
        self.thread.start()
        
        # --- UI Construction ---
        self.setup_grid()
        self.create_header_bar()
        self.create_launch_controls()
        self.create_main_container()
        self.create_activity_log()
        
        # Start logging check
        self.after(100, self.process_logs)
        
        # Handle Close
        self.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        self.log("System Initialized. Ready for Phase 1.")

    def setup_grid(self):
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(0, weight=0)  # Header Bar
        self.grid_rowconfigure(1, weight=0)  # Window Management
        self.grid_rowconfigure(2, weight=1)  # Main Content (Sidebar + Content + Sessions)
        self.grid_rowconfigure(3, weight=0)  # Activity Logs

    def create_header_bar(self):
        """Fluid header with proportional spacing."""
        header = ctk.CTkFrame(self, fg_color=COLOR["surface"], corner_radius=0, height=50)
        header.grid(row=0, column=0, sticky="nsew")
        header.grid_columnconfigure(0, weight=1) # Logo area
        header.grid_columnconfigure(1, weight=1) # Profile area

        # Logo area (Left)
        logo_frame = ctk.CTkFrame(header, fg_color="transparent")
        logo_frame.grid(row=0, column=0, padx=20, sticky="w")
        
        try:
            if os.path.exists(self.logo_path):
                logo_img = ctk.CTkImage(light_image=PILImage.open(self.logo_path), dark_image=PILImage.open(self.logo_path), size=(32, 32))
                ctk.CTkLabel(logo_frame, text="", image=logo_img).pack(side="left", padx=(0, 10))
        except:
            pass
            
        ctk.CTkLabel(logo_frame, text="ROX-", font=("Inter", 20, "bold"), text_color=COLOR["text"]).pack(side="left")
        ctk.CTkLabel(logo_frame, text="SHOOTER", font=("Inter", 20, "bold"), text_color=COLOR["primary"]).pack(side="left")
        ctk.CTkLabel(logo_frame, text=" v1.1", font=("Inter", 12), text_color=COLOR["text_sec"]).pack(side="left", padx=(4, 0))

        # User Profile (Right)
        avatar_frame = ctk.CTkFrame(header, fg_color="transparent")
        avatar_frame.grid(row=0, column=1, padx=20, sticky="e")
        
        initial = self.user_id[0].upper() if self.user_id else "U"
        avatar = ctk.CTkButton(
            avatar_frame, text=initial, width=32, height=32, corner_radius=16,
            fg_color=COLOR["primary"], hover=False,
            text_color="#FFFFFF", font=("Inter", 14, "bold"), state="normal"
        )
        avatar.pack(side="right", padx=(8, 0))
        ctk.CTkLabel(avatar_frame, text=self.user_id, font=("Inter", 12, "bold"), text_color=COLOR["text"]).pack(side="right")

    def create_launch_controls(self):
        # ── Window Management Bar ──
        self.launch_frame = ctk.CTkFrame(self, fg_color=COLOR["surface"], corner_radius=12, border_width=1, border_color=COLOR["border"])
        self.launch_frame.grid(row=1, column=0, padx=24, pady=(8, 4), sticky="nsew")
        self.launch_frame.grid_columnconfigure(0, weight=1)
        self.launch_frame.grid_columnconfigure(1, weight=1)

        # Left group
        left_group = ctk.CTkFrame(self.launch_frame, fg_color="transparent")
        left_group.grid(row=0, column=0, padx=15, pady=10, sticky="w")

        ctk.CTkLabel(left_group, text="🖥  Control", font=("Inter", 12, "bold"), text_color=COLOR["text"]).pack(side="left", padx=(0, 10))
        
        self.entry_num_windows = ctk.CTkEntry(left_group, width=44, height=30, corner_radius=6,
                                               fg_color=COLOR["input_bg"], border_color=COLOR["input_border"],
                                               text_color=COLOR["text"], placeholder_text="1", font=("Inter", 12))
        self.entry_num_windows.insert(0, "1")
        self.entry_num_windows.pack(side="left", padx=(0, 8))

        self.btn_launch = ctk.CTkButton(left_group, text="Initialize", height=30, width=100, corner_radius=6,
                                         fg_color=COLOR["primary"], hover_color=COLOR["primary_hov"],
                                         text_color="#FFFFFF", font=("Inter", 11, "bold"), command=self.on_launch_windows)
        self.btn_launch.pack(side="left", padx=(0, 6))

        self.btn_terminate = ctk.CTkButton(left_group, text="Close All", height=30, width=90, corner_radius=6,
                                            fg_color=COLOR["danger"], hover_color=COLOR["danger_hov"],
                                            text_color="#FFFFFF", font=("Inter", 11, "bold"), command=self.on_terminate_all)
        self.btn_terminate.pack(side="left")

        # Right group
        right_group = ctk.CTkFrame(self.launch_frame, fg_color="transparent")
        right_group.grid(row=0, column=1, padx=15, pady=10, sticky="e")

        self.btn_reset = ctk.CTkButton(right_group, text="System Reset", width=110, height=30, corner_radius=6,
                                        fg_color=COLOR["surface_alt"], hover_color=COLOR["border"],
                                        text_color=COLOR["text"], border_width=1, border_color=COLOR["border"],
                                        font=("Inter", 11, "bold"), command=self.on_reset)
        self.btn_reset.pack(side="right", padx=(6, 0))

        self.btn_clear_log = ctk.CTkButton(right_group, text="Clear Logs", width=90, height=30, corner_radius=6,
                                            fg_color=COLOR["surface_alt"], hover_color=COLOR["border"],
                                            text_color=COLOR["text"], border_width=1, border_color=COLOR["border"],
                                            font=("Inter", 11, "bold"), command=self.on_clear_log)
        self.btn_clear_log.pack(side="right")

    def create_main_container(self):
        """Main layout: Sidebar (Left), Content (Middle), Active Sessions (Right)"""
        self.main_container = ctk.CTkFrame(self, fg_color="transparent")
        self.main_container.grid(row=2, column=0, padx=24, pady=8, sticky="nsew")

        # Layout: Sidebar (0), Content (1), Sessions (2)
        self.main_container.grid_columnconfigure(0, weight=0) # Sidebar
        self.main_container.grid_columnconfigure(1, weight=10) # Content area (Takes most space)
        self.main_container.grid_columnconfigure(2, weight=4) # Active sessions
        self.main_container.grid_rowconfigure(0, weight=1)

        # ── Sidebar Navigation ──
        self.sidebar_frame = ctk.CTkScrollableFrame(
            self.main_container, width=180, corner_radius=12,
            fg_color=COLOR["surface"], border_width=1, border_color=COLOR["border"],
            label_text="", scrollbar_button_color=COLOR["surface_alt"]
        )
        self.sidebar_frame.grid(row=0, column=0, padx=(0, 16), sticky="nsew")
        # Removed grid_propagate(False) as it's not supported/needed for CTkScrollableFrame

        ctk.CTkLabel(self.sidebar_frame, text=" NAVIGATION", font=("Inter", 11, "bold"), text_color=COLOR["text_sec"]).pack(anchor="w", padx=10, pady=(10, 12))

        # Content area (to hold different tab frames)
        self.content_area = ctk.CTkFrame(
            self.main_container, corner_radius=12,
            fg_color=COLOR["surface"], border_width=1, border_color=COLOR["border"]
        )
        self.content_area.grid(row=0, column=1, padx=(0, 16), sticky="nsew")

        # ── Active Sessions Panel (Right Side) ──
        self.active_windows_frame = ctk.CTkScrollableFrame(
            self.main_container, label_text="⚡ Sessions Monitor",
            label_font=("Inter", 14, "bold"),
            fg_color=COLOR["surface"], corner_radius=12,
            border_width=1, border_color=COLOR["border"],
            label_fg_color=COLOR["surface"]
        )
        self.active_windows_frame.grid(row=0, column=2, sticky="nsew")

        # Internal state for tab switching
        self.nav_buttons = {}
        self.tab_frames = {}
        self.current_tab = None

        # Define tabs and their icons
        tabs = [
            ("Target Contacts", "👥"),
            ("Subject & Body", "✉️"),
            ("Content Options", "📑"),
            ("TFN Number", "📞"),
            ("Advanced Params", "⚙️"),
            ("Blasting Center", "🚀"),
            ("System Tags", "🏷️")
        ]

        for name, icon in tabs:
            # IMPORTANT: All tab content is now in CTkScrollableFrame for full responsiveness
            frame = ctk.CTkScrollableFrame(self.content_area, fg_color="transparent", corner_radius=0)
            self.tab_frames[name] = frame
            
            # Create nav button
            btn = ctk.CTkButton(
                self.sidebar_frame, text=f" {icon}  {name}", 
                anchor="w", height=40, corner_radius=8,
                fg_color="transparent", text_color=COLOR["text"],
                hover_color=COLOR["surface_alt"], font=("Inter", 12, "bold"),
                command=lambda n=name: self.show_tab(n)
            )
            btn.pack(fill="x", padx=10, pady=2)
            self.nav_buttons[name] = btn

        # Define specific frame variables for backward compatibility
        self.tab_data = self.tab_frames["Target Contacts"]
        self.tab_subject_body = self.tab_frames["Subject & Body"]
        self.tab_content = self.tab_frames["Content Options"]
        self.tab_tfn = self.tab_frames["TFN Number"]
        self.tab_settings = self.tab_frames["Advanced Params"]
        self.tab_blasting = self.tab_frames["Blasting Center"]
        self.tab_tags = self.tab_frames["System Tags"]

        self.setup_tabs()
        self.show_tab("Target Contacts")

    def show_tab(self, name):
        """Switches the visible content frame and updates sidebar styling."""
        if self.current_tab:
            self.tab_frames[self.current_tab].pack_forget()
            self.nav_buttons[self.current_tab].configure(fg_color="transparent", text_color=COLOR["text"])

        self.tab_frames[name].pack(fill="both", expand=True)
        self.nav_buttons[name].configure(fg_color=COLOR["primary"], text_color="#FFFFFF")
        self.current_tab = name

    def setup_tabs(self):
        # --- Tab 1: Data ---
        self.data_frame = ctk.CTkFrame(self.tab_data, fg_color="transparent")
        self.data_frame.pack(fill="both", expand=True, padx=24, pady=20)

        ctk.CTkLabel(self.data_frame, text="Target Contacts", font=("Inter", 14, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 8))

        btn_row = ctk.CTkFrame(self.data_frame, fg_color="transparent")
        btn_row.pack(fill="x", pady=(0, 10))

        self.btn_load_data = ctk.CTkButton(
            btn_row, text="📂  Load CSV / Excel / TXT", height=36, corner_radius=8,
            fg_color=COLOR["info"], hover_color=COLOR["info_hov"],
            text_color="#FFFFFF", font=("Inter", 12, "bold"),
            command=self.handle_load_file
        )
        self.btn_load_data.pack(side="left")

        self.text_emails = ctk.CTkTextbox(
            self.data_frame, height=300, font=("Consolas", 12),
            fg_color=COLOR["input_bg"], text_color=COLOR["text"],
            border_width=1, border_color=COLOR["input_border"], corner_radius=8
        )
        self.text_emails.pack(fill="both", expand=True, pady=(0, 10))
        self.text_emails.bind("<KeyRelease>", self.update_email_counter)

        self.lbl_email_counter = ctk.CTkLabel(
            self.data_frame, text="Total Emails: 0",
            font=("Inter", 14, "bold"), text_color=COLOR["primary"]
        )
        self.lbl_email_counter.pack(pady=(0, 20))

        # --- Tab 2: Subject & Body ---
        sb_frame = ctk.CTkFrame(self.tab_subject_body, fg_color="transparent")
        sb_frame.pack(fill="both", expand=True, padx=24, pady=20)

        ctk.CTkLabel(sb_frame, text="Blasting Subject", font=("Inter", 14, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 6))
        self.entry_subject = ctk.CTkEntry(
            sb_frame, placeholder_text="e.g. $invoice_no7 — Important Notice", font=("Inter", 13), height=38,
            corner_radius=8, fg_color=COLOR["input_bg"], border_color=COLOR["input_border"],
            text_color=COLOR["text"]
        )
        self.entry_subject.pack(fill="x", pady=(0, 16))

        ctk.CTkLabel(sb_frame, text="Main Message Body", font=("Inter", 14, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 6))
        self.text_body = ctk.CTkTextbox(
            sb_frame, height=300, font=("Inter", 13),
            fg_color=COLOR["input_bg"], text_color=COLOR["text"],
            border_width=1, border_color=COLOR["input_border"], corner_radius=8
        )
        self.text_body.pack(fill="both", expand=True, pady=(0, 20))
        
        # --- Tab 3: Content ---
        # Content frame as a container inside the already scrollable tab_content
        self.content_frame = ctk.CTkFrame(self.tab_content, fg_color="transparent")
        self.content_frame.pack(fill="both", expand=True, padx=24, pady=20)

        # Row 1: Conversion Type
        self.conv_group = ctk.CTkFrame(self.content_frame, fg_color="transparent")
        self.conv_group.pack(fill="x", pady=(0, 24))

        ctk.CTkLabel(self.conv_group, text="Conversion Type", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 4))
        self.conv_options = [
            "None", "Standard PDF", "PNG Image", "Secure PDF",
            "Secure PPTX", "Secure Excel", "Standard PPTX", "Standard Excel"
        ]
        self.dropdown_conversion = ctk.CTkComboBox(
            self.conv_group, values=self.conv_options, width=320, height=36, corner_radius=8,
            fg_color=COLOR["input_bg"], border_color=COLOR["input_border"],
            text_color=COLOR["text"], button_color=COLOR["primary"],
            button_hover_color=COLOR["primary_hov"], dropdown_fg_color=COLOR["surface"],
            dropdown_text_color=COLOR["text"], dropdown_hover_color=COLOR["surface_alt"]
        )
        self.dropdown_conversion.set("Standard PDF")
        self.dropdown_conversion.pack(anchor="w", pady=(0, 10))
        
        self.btn_preview = ctk.CTkButton(
            self.conv_group, text="Preview Selected", height=32, corner_radius=8,
            fg_color=COLOR["accent"], hover_color=COLOR["accent_hov"],
            text_color="#FFFFFF", font=("Inter", 11, "bold"), command=self.on_preview_attachment
        )
        self.btn_preview.pack(anchor="w")

        # Row 2: Filename Mode
        self.file_group = ctk.CTkFrame(self.content_frame, fg_color="transparent")
        self.file_group.pack(fill="x", pady=(0, 24))

        ctk.CTkLabel(self.file_group, text="Filename Mode", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 4))
        self.file_mode = ctk.CTkSegmentedButton(
            self.file_group, values=["Random", "Custom"], command=self.toggle_filename_entry,
            selected_color=COLOR["primary"], selected_hover_color=COLOR["primary_hov"],
            unselected_color=COLOR["surface_alt"], unselected_hover_color=COLOR["border"],
            text_color=COLOR["text"], text_color_disabled=COLOR["muted"],
            font=("Inter", 12, "bold"), corner_radius=8
        )
        self.file_mode.set("Random")
        self.file_mode.pack(anchor="w")

        # Row 3: Custom Filename Entry (Packed inside file_group to appear below button)
        self.filename_frame = ctk.CTkFrame(self.file_group, fg_color="transparent")
        ctk.CTkLabel(self.filename_frame, text="Custom Filename (Supports Tags):", font=("Inter", 12), text_color=COLOR["text_sec"]).pack(anchor="w", pady=(10, 0))
        self.entry_filename = ctk.CTkEntry(
            self.filename_frame, placeholder_text="e.g. Invoice_$invoice_no7", width=400, height=36,
            corner_radius=8, fg_color=COLOR["input_bg"], border_color=COLOR["input_border"],
            text_color=COLOR["text"]
        )
        self.entry_filename.pack(anchor="w", pady=5)
        # Hidden by default, toggle handles packing


        # Row 3: HTML Content
        ctk.CTkLabel(self.content_frame, text="HTML Content", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(4, 6))
        self.text_html = ctk.CTkTextbox(
            self.content_frame, font=("Consolas", 12), height=400,
            fg_color=COLOR["input_bg"], text_color=COLOR["text"],
            border_width=1, border_color=COLOR["input_border"], corner_radius=8
        )
        self.text_html.pack(fill="x", pady=(0, 40))

        
        # --- Tab 4: TFN ---
        tfn_frame = ctk.CTkFrame(self.tab_tfn, fg_color="transparent")
        tfn_frame.pack(fill="both", expand=True, padx=24, pady=20)

        ctk.CTkLabel(tfn_frame, text="Toll-Free Number (TFN)", font=("Inter", 16, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 16))
        
        tfn_card = ctk.CTkFrame(tfn_frame, fg_color=COLOR["surface_alt"], corner_radius=10, border_width=1, border_color=COLOR["border"])
        tfn_card.pack(fill="x", pady=(0, 8))
        tfn_inner = ctk.CTkFrame(tfn_card, fg_color="transparent")
        tfn_inner.pack(fill="x", padx=20, pady=16)
        ctk.CTkLabel(tfn_inner, text="Configure TFN for Replacement", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 6))
        self.entry_tfn = ctk.CTkEntry(
            tfn_inner, placeholder_text="e.g. +1-800-XXX-XXXX", width=400, height=36,
            corner_radius=8, fg_color=COLOR["input_bg"], border_color=COLOR["input_border"], text_color=COLOR["text"]
        )
        self.entry_tfn.pack(anchor="w", pady=(0, 20))

        # --- Tab 5: Settings ---
        settings_frame = ctk.CTkFrame(self.tab_settings, fg_color="transparent")
        settings_frame.pack(fill="both", expand=True, padx=24, pady=20)

        ctk.CTkLabel(settings_frame, text="Advanced Parameters", font=("Inter", 16, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 16))

        ctk.CTkLabel(settings_frame, text="⚠  Tuning these delays affects speed and stability. Lower values may cause UI glitches.", font=("Inter", 11), text_color=COLOR["muted"]).pack(anchor="w", pady=(0, 16))

        # Presets UI
        preset_frame = ctk.CTkFrame(settings_frame, fg_color="transparent")
        preset_frame.pack(fill="x", pady=(0, 16))
        ctk.CTkLabel(preset_frame, text="Speed Presets:", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(side="left", padx=(0, 16))
        self.preset_selector = ctk.CTkSegmentedButton(
            preset_frame, values=["Normal (Default)", "Medium (Faster)", "Fast (Risky)"],
            command=self.apply_delay_preset,
            selected_color=COLOR["primary"], selected_hover_color=COLOR["primary_hov"],
            unselected_color=COLOR["surface_alt"], unselected_hover_color=COLOR["border"],
            font=("Inter", 12, "bold")
        )
        self.preset_selector.set("Normal (Default)")
        self.preset_selector.pack(side="left")

        self.delay_entries = {}
        
        delay_configs = [
            ("General Delays", [
                ("interval", "Delay Between Emails", "Wait time between sending emails", "2.0")
            ]),
            ("Page & Network Delays", [
                ("gmail_load", "Gmail Page Load", "Wait for Gmail to load initially", "3.0"),
                ("upload", "Attachment Upload", "Wait for file upload to complete", "3.5"),
                ("inline_image", "Inline Image Process", "Wait for pasted HTML images to process", "2.5")
            ]),
            ("Automation Action Delays", [
                ("compose", "Compose Click", "Wait after clicking Compose button", "2.25"),
                ("discard", "Discard Drafts", "Wait after cleaning up old drafts", "1.0"),
                ("post_send", "Post-Send Wait", "Wait to confirm email was sent", "2.0"),
                ("retry", "Conversion Retry", "Wait between conversion retries (PDF/Image)", "1.5")
            ]),
            ("Micro Delays", [
                ("micro", "Micro Action Delay", "Delay for typing, clicking, and Enter/Tab", "0.5")
            ])
        ]

        for group_name, configs in delay_configs:
            group_lbl = ctk.CTkLabel(settings_frame, text=group_name, font=("Inter", 14, "bold"), text_color=COLOR["primary"])
            group_lbl.pack(anchor="w", pady=(10, 4))
            
            for key, title, desc, default in configs:
                card = ctk.CTkFrame(settings_frame, fg_color=COLOR["surface_alt"], corner_radius=10, border_width=1, border_color=COLOR["border"])
                card.pack(fill="x", pady=(0, 8))
                
                inner = ctk.CTkFrame(card, fg_color="transparent")
                inner.pack(fill="x", padx=16, pady=12)
                
                text_frame = ctk.CTkFrame(inner, fg_color="transparent")
                text_frame.pack(side="left", fill="x", expand=True)
                
                ctk.CTkLabel(text_frame, text=title, font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(anchor="w")
                ctk.CTkLabel(text_frame, text=desc, font=("Inter", 11), text_color=COLOR["text_sec"]).pack(anchor="w")
                
                entry = ctk.CTkEntry(
                    inner, width=70, height=32, corner_radius=6,
                    fg_color=COLOR["input_bg"], border_color=COLOR["input_border"], text_color=COLOR["text"], font=("Inter", 12)
                )
                entry.insert(0, default)
                entry.pack(side="right")
                self.delay_entries[key] = entry
                
                if key == "interval":
                    self.entry_delay = entry

        # --- Tab 6: Blasting ---
        # Tab content is already scrollable.
        self.tab_blasting.grid_columnconfigure(0, weight=1)

        # Center card container
        blasting_card = ctk.CTkFrame(
            self.tab_blasting, 
            fg_color=COLOR["surface_alt"], 
            corner_radius=20, 
            border_width=1, 
            border_color=COLOR["border"]
        )
        blasting_card.pack(pady=30, padx=20, anchor="center")

        # Decorative and descriptive labels
        ctk.CTkLabel(blasting_card, text="🚀", font=("Inter", 36)).pack(pady=(24, 8))
        ctk.CTkLabel(blasting_card, text="Blasting Center", font=("Inter", 18, "bold"), text_color=COLOR["text"]).pack(pady=(0, 6))
        ctk.CTkLabel(blasting_card, text="Ready to start the automated mailing engine?\nEnsure all parameters are configured.", 
                     font=("Inter", 12), text_color=COLOR["text_sec"], justify="center").pack(pady=(0, 24), padx=30)

        self.btn_start_blasting = ctk.CTkButton(
            blasting_card, text="START BLASTING", height=48, width=240,
            corner_radius=24, font=("Inter", 14, "bold"),
            fg_color=COLOR["primary"], hover_color=COLOR["primary_hov"],
            text_color="#FFFFFF", command=self.on_start_blasting
        )
        self.btn_start_blasting.pack(pady=(0, 36), padx=40)

        # --- Tab 7: Tags ---
        self.setup_tags_tab()

    def apply_delay_preset(self, preset_name):
        presets = {
            "Normal (Default)": {
                "interval": "2.0", "gmail_load": "3.0", "upload": "3.5",
                "inline_image": "2.5", "compose": "2.25", "discard": "1.0",
                "post_send": "2.0", "retry": "1.5", "micro": "0.5"
            },
            "Medium (Faster)": {
                "interval": "1.0", "gmail_load": "2.0", "upload": "2.5",
                "inline_image": "1.5", "compose": "1.5", "discard": "0.5",
                "post_send": "1.0", "retry": "1.0", "micro": "0.3"
            },
            "Fast (Risky)": {
                "interval": "0.0", "gmail_load": "1.5", "upload": "1.5",
                "inline_image": "1.0", "compose": "0.5", "discard": "0.2",
                "post_send": "0.5", "retry": "0.5", "micro": "0.1"
            }
        }
        if preset_name in presets and hasattr(self, 'delay_entries'):
            vals = presets[preset_name]
            for k, v in vals.items():
                if k in self.delay_entries:
                    self.delay_entries[k].delete(0, "end")
                    self.delay_entries[k].insert(0, v)

    def setup_tags_tab(self):
        container = ctk.CTkFrame(self.tab_tags, fg_color="transparent")
        container.pack(fill="both", expand=True, padx=24, pady=20)

        ctk.CTkLabel(container, text="Available System Tags", font=("Inter", 16, "bold"), text_color=COLOR["text"]).pack(anchor="w", pady=(0, 4))
        ctk.CTkLabel(container, text="Use these tags in Subject, Body, HTML Content, or Custom Filenames.\nClick 'Copy' to quickly add a tag to your clipboard.",
                      font=("Inter", 12), justify="left", wraplength=550, text_color=COLOR["text_sec"]).pack(anchor="w", pady=(0, 16))

        tags_info = [
            ("$email", "Recipient's full email address.\nExample: god.mailer@gmail.com", COLOR["info"]),
            ("$n_mail", "Email username (portion before '@').\nExample: god.mailer", COLOR["accent"]),
            ("$invoice_no", "Random alphanumeric invoice number.\nExample: 22250174NCN8", COLOR["success"]),
            ("$number6", "Random 6-digit number.\nExample: 767059", COLOR["warn"]),
            ("$random8", "8-character alphanumeric string.\nExample: G33ACOH3", "#EC4899"),
            ("$tfn", "Toll-free number from designated field.\nExample: 1-800-555-0199", COLOR["primary"]),
            ("$word3", "Random 3-letter uppercase word.\nExample: TKS", COLOR["info"]),
            ("$word4", "Random 4-letter uppercase word.\nExample: MJJR", COLOR["accent"]),
            ("$date", "Randomly formatted date (DD-MM-YY).\nExample: 11-05-26", COLOR["muted"]),
            ("$digi_no", "Long token identifier (UUID).\nExample: badb750c-0a97-4229...", COLOR["success"]),
            ("$usa_name", "Standard full name (First Last).\nExample: Garrett Benton", COLOR["warn"]),
            ("$usa_rnd", "Randomized name format.\nExample: J.D. Smith or Smith, John M.", "#EC4899"),
            ("$str_ing", "Random string (5 or 11 letters).\nExample: QuWOmbxzqmL", COLOR["primary"]),
            ("$mail_rnd", "Username + 8 random digits.\nExample: god.mailer_41091156", COLOR["info"]),
            ("$fn_usa", "Outputs only the first name.\nExample: Misty", COLOR["accent"]),
            ("$charge", "Random charge amount ($10-$99).\nExample: 59", COLOR["success"]),
        ]

        for tag, desc, color in tags_info:
            frame = ctk.CTkFrame(container, fg_color=COLOR["tag_card"], corner_radius=10, border_width=1, border_color=COLOR["border"])
            frame.pack(fill="x", pady=4)

            lbl_tag = ctk.CTkLabel(frame, text=tag, font=("Consolas", 13, "bold"), text_color=color, width=140, anchor="w")
            lbl_tag.pack(side="left", padx=16, pady=12)

            lbl_desc = ctk.CTkLabel(frame, text=desc, font=("Inter", 12), justify="left", text_color=COLOR["text"])
            lbl_desc.pack(side="left", padx=10, pady=12, fill="x", expand=True)

            btn_copy = ctk.CTkButton(
                frame, text="Copy", width=60, height=28, corner_radius=6,
                fg_color=COLOR["surface"], hover_color=COLOR["border"],
                text_color=COLOR["primary"], border_width=1, border_color=COLOR["primary"],
                font=("Inter", 11, "bold"),
                command=lambda t=tag: self.copy_to_clipboard(t)
            )
            btn_copy.pack(side="right", padx=16)

        ctk.CTkLabel(container, text="Note: (length) is optional and defaults to 6 if not specified.", font=("Inter", 11, "italic"), text_color=COLOR["muted"]).pack(anchor="w", pady=(10, 30))

    def copy_to_clipboard(self, text):
        self.clipboard_clear()
        self.clipboard_append(text)


    def create_activity_log(self):
        # ── Activity Logs ──
        self.log_frame = ctk.CTkFrame(self, fg_color=COLOR["surface"], corner_radius=0, border_width=1, border_color=COLOR["border"])
        self.log_frame.grid(row=3, column=0, sticky="nsew")
        
        log_header = ctk.CTkFrame(self.log_frame, fg_color="transparent")
        log_header.pack(fill="x", padx=16, pady=(12, 0))
        ctk.CTkLabel(log_header, text="✨  Activity Logs", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(side="left")

        self.log_textbox = ctk.CTkTextbox(
            self.log_frame, font=("Consolas", 11.5), state="disabled",
            fg_color=COLOR["log_bg"], text_color=COLOR["log_text"],
            corner_radius=8, border_width=1, border_color=COLOR["border"],
            height=140
        )
        self.log_textbox.pack(fill="both", expand=True, padx=16, pady=(8, 14))

    # --- UI Logic & Handlers ---
    def log(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_queue.put(f"[{timestamp}] {message}\n")

    def process_logs(self):
        if not self.winfo_exists():
            return
            
        while not self.log_queue.empty():
            try:
                msg = self.log_queue.get_nowait()
                if self.log_textbox.winfo_exists():
                    self.log_textbox.configure(state="normal")
                    self.log_textbox.insert("end", msg)
                    self.log_textbox.see("end")
                    self.log_textbox.configure(state="disabled")
            except:
                break
                
        if self.winfo_exists():
            self.after(100, self.process_logs)

    def update_email_counter(self, event=None):
        content = self.text_emails.get("1.0", "end-1c").strip()
        emails = [e for e in content.split("\n") if e.strip()]
        self.lbl_email_counter.configure(text=f"Total Emails: {len(emails)}")

    def on_launch_windows(self):
        try:
            count_str = self.entry_num_windows.get()
            count = int(count_str) if count_str.isdigit() else 1
            self.run_coro(self.launch_and_arrange(count))
        except Exception as e:
            self.log(f"Launch Error: {e}")

    async def launch_and_arrange(self, count):
        """Launch all sessions sequentially, then tile them split-screen."""
        for i in range(1, count + 1):
            if i not in self.contexts:
                await self.launch_browser_task(i)

        # Silently arrange all active windows in split-screen
        await self.arrange_windows_split()

    async def arrange_windows_split(self):
        """Tile all active browser windows in a grid across the screen using CDP."""
        import math
        active_ids = sorted(self.contexts.keys())
        n = len(active_ids)
        if n == 0:
            return

        screen_w = self.winfo_screenwidth()
        screen_h = self.winfo_screenheight()

        cols = math.ceil(math.sqrt(n))
        rows = math.ceil(n / cols)
        win_w = screen_w // cols
        win_h = screen_h // rows

        for idx, wid in enumerate(active_ids):
            x = (idx % cols) * win_w
            y = (idx // cols) * win_h
            try:
                page = self.contexts[wid]["page"]
                cdp = await page.context.new_cdp_session(page)
                win_info = await cdp.send("Browser.getWindowForTarget")
                await cdp.send("Browser.setWindowBounds", {
                    "windowId": win_info["windowId"],
                    "bounds": {"left": x, "top": y, "width": win_w, "height": win_h, "windowState": "normal"}
                })
            except:
                pass

    async def launch_browser_task(self, window_id):
        try:
            self.log(f"Initializing Session {window_id}...")
            profile_path = os.path.abspath(f"profiles/profile_{window_id}")
            os.makedirs(profile_path, exist_ok=True)
            
            # Force remove lock file if it exists (prevents ProcessSingleton errors)
            lock_file = os.path.join(profile_path, "SingletonLock")
            if os.path.exists(lock_file):
                try:
                    os.remove(lock_file)
                except:
                    pass
            
            # Start Playwright for this context
            pw = await async_playwright().start()
            context = await pw.chromium.launch_persistent_context(
                user_data_dir=profile_path,
                channel="chrome",
                headless=False,
                args=["--incognito", "--disable-blink-features=AutomationControlled"]
            )
            
            # Initial page: reuse default tab if exists, else create new
            if context.pages:
                page = context.pages[0]
            else:
                page = await context.new_page()
            
            # Store in app state
            self.contexts[window_id] = {
                "pw": pw,
                "context": context,
                "page": page
            }
            
            # Update UI (thread-safe via after)
            self.after(0, lambda: self.add_window_row(window_id))
            self.log(f"Session {window_id} ready.")
            
            # Automatically open Gmail on initialization
            await self.open_gmail_task(window_id)
            
        except Exception as e:
            self.log(f"Session {window_id} failed: {e}")

    def add_window_row(self, window_id):
        # Create a styled row card in the scrollable frame
        row = ctk.CTkFrame(self.active_windows_frame, fg_color=COLOR["surface_alt"], corner_radius=12, border_width=1, border_color=COLOR["border"])
        row.pack(fill="x", padx=10, pady=6)

        # Header with session ID
        ctk.CTkLabel(row, text=f"⚡ Session {window_id}", font=("Inter", 13, "bold"), text_color=COLOR["text"]).pack(side="top", anchor="w", padx=15, pady=(12, 8))

        # Button row below label
        btn_row = ctk.CTkFrame(row, fg_color="transparent")
        btn_row.pack(fill="x", padx=12, pady=(0, 12))

        btn_gmail = ctk.CTkButton(
            btn_row, text="Open Browser", height=32, corner_radius=8,
            fg_color=COLOR["info"], hover_color=COLOR["info_hov"], text_color="#FFFFFF",
            font=("Inter", 11, "bold"),
            command=lambda: self.run_coro(self.open_gmail_task(window_id))
        )
        btn_gmail.pack(side="left", fill="x", expand=True, padx=(0, 6))

        btn_close = ctk.CTkButton(
            btn_row, text="End", width=60, height=32, corner_radius=8,
            fg_color=COLOR["surface"], hover_color=COLOR["border"], text_color=COLOR["danger"],
            border_width=1, border_color=COLOR["danger"],
            font=("Inter", 11, "bold"),
            command=lambda: self.run_coro(self.close_window_task(window_id))
        )
        btn_close.pack(side="right")

        # Save the row frame for deletion later
        self.contexts[window_id]["ui_row"] = row

    async def open_gmail_task(self, window_id):
        if window_id in self.contexts:
            page = self.contexts[window_id]["page"]
            self.log(f"Navigating Window {window_id} to Gmail...")
            await page.goto("https://mail.google.com")
            await page.bring_to_front()

    async def close_window_task(self, window_id):
        if window_id in self.contexts:
            # Get reference but don't pop yet to keep tracking
            data = self.contexts.get(window_id)
            if not data: return
            
            self.log(f"Closing Window {window_id}...")
            try:
                if "page" in data:
                    await data["page"].close()
                if "context" in data:
                    await data["context"].close()
                if "pw" in data:
                    await data["pw"].stop()
                self.log(f"Window {window_id} closed successfully.")
            except Exception as e:
                self.log(f"Error closing Window {window_id}: {e}")
            finally:
                # Remove from tracking after attempt
                if window_id in self.contexts:
                    self.contexts.pop(window_id)
            
            # Update UI
            if "ui_row" in data:
                try:
                    self.after(0, data["ui_row"].destroy)
                except:
                    pass

    def on_terminate_all(self):
        self.log("Terminating all windows...")
        ids = list(self.contexts.keys())
        for window_id in ids:
            self.run_coro(self.close_window_task(window_id))

    def on_clear_log(self):
        self.log_textbox.configure(state="normal")
        self.log_textbox.delete("1.0", "end")
        self.log_textbox.configure(state="disabled")

    def on_reset(self):
        self.entry_num_windows.delete(0, "end")
        self.entry_num_windows.insert(0, "1")
        self.text_emails.delete("1.0", "end")
        self.entry_subject.delete(0, "end")
        self.text_body.delete("1.0", "end")
        self.text_html.delete("1.0", "end")
        self.entry_tfn.delete(0, "end")
        self.update_email_counter()
        self.log("Application state reset.")

    def handle_load_file(self):
        from tkinter import filedialog
        import pandas as pd
        import re
        
        file_path = filedialog.askopenfilename(filetypes=[("Data Files", "*.csv *.xlsx *.xls *.txt")])
        if not file_path:
            return
            
        try:
            self.log(f"Loading data from {os.path.basename(file_path)}...")
            if file_path.endswith(".csv"):
                df = pd.read_csv(file_path)
                text = df.to_string()
            elif file_path.endswith((".xlsx", ".xls")):
                df = pd.read_excel(file_path)
                text = df.to_string()
            else:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            
            # Extract emails using regex
            emails = re.findall(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', text)
            unique_emails = list(dict.fromkeys(emails)) # Remove duplicates while preserving order
            
            self.text_emails.delete("1.0", "end")
            self.text_emails.insert("1.0", "\n".join(unique_emails))
            self.update_email_counter()
            self.log(f"Successfully imported {len(unique_emails)} unique emails.")
            
        except Exception as e:
            self.log(f"Error loading file: {e}")

    def toggle_filename_entry(self, value):
        if value == "Custom":
            self.filename_frame.pack(fill="x", pady=(10, 0))
        else:
            self.filename_frame.pack_forget()

    def on_preview_attachment(self):
        self.run_coro(self.preview_attachment_task())

    def on_preview_specific(self, format_name):
        self.run_coro(self.preview_attachment_task(format_override=format_name))

    async def preview_attachment_task(self, format_override=None):
        html_template = self.text_html.get("1.0", "end-1c")
        if not html_template.strip():
            self.log("Error: No HTML content to preview.")
            return

        conversion_type = format_override if format_override else self.dropdown_conversion.get()
        if conversion_type == "None":
            self.log("Error: Conversion type is 'None'. Choose a type to preview.")
            return

        self.log(f"Generating preview for {conversion_type}...")
        
        # Parse with dummy data
        tfn = self.entry_tfn.get() or "1-800-PREVIEW"
        parsed_html = self.parser.parse(html_template, "preview@example.com", tfn)
        
        # Extension
        ext = ".pdf"
        if "Image" in conversion_type and "PDF" not in conversion_type and "PPTX" not in conversion_type and "XLS" not in conversion_type and "PPT" not in conversion_type and "Excel" not in conversion_type:
            ext = ".png"
        elif "HTML Mosaic" in conversion_type or "CSS Scrambled" in conversion_type or "BiDi Email" in conversion_type:
            ext = ".html"
        elif "PPT" in conversion_type:
            ext = ".pptx"
        elif "XLS" in conversion_type or "Excel" in conversion_type:
            ext = ".xlsx"
            
        filename = f"preview_test{ext}"
        path = None
        
        try:
            if conversion_type == "Standard PDF":
                path = await self.converter.html_to_pdf(parsed_html, filename)
            elif conversion_type == "PNG Image":
                path = await self.converter.html_to_image(parsed_html, filename)
            elif conversion_type == "Secure PDF":
                path = await self.converter.html_to_image_pdf(parsed_html, filename)
            elif conversion_type == "Secure PPTX":
                path = await self.converter.html_to_image_pptx(parsed_html, filename)
            elif conversion_type == "Secure Excel":
                path = await self.converter.html_to_image_xls(parsed_html, filename)
            elif conversion_type == "Standard PPTX":
                path = await self.converter.html_to_pdf_pptx(parsed_html, filename)
            elif conversion_type == "Standard Excel":
                path = await self.converter.html_to_pdf_xls(parsed_html, filename)
            
            if path and os.path.exists(path):
                self.log(f"Preview generated: {path}")
                # Open the file based on OS
                if sys.platform == "darwin":
                    subprocess.run(["open", path])
                elif sys.platform == "win32":
                    os.startfile(path)
                else:
                    try:
                        subprocess.run(["xdg-open", path])
                    except:
                        self.log(f"File saved at: {path} (Could not open automatically)")
            else:
                self.log("Error: Failed to generate preview file.")
        except Exception as e:
            self.log(f"Preview Error: {e}")

    def on_preview_all(self):
        self.run_coro(self.preview_all_task())

    async def preview_all_task(self):
        html_template = self.text_html.get("1.0", "end-1c")
        if not html_template.strip():
            self.log("Error: No HTML content to preview.")
            return

        self.log("Generating previews for ALL formats. This may take a moment...")
        
        # Parse with dummy data
        tfn = self.entry_tfn.get() or "1-800-PREVIEW"
        parsed_html = self.parser.parse(html_template, "preview@example.com", tfn)
        
        tasks = [
            ("PDF", self.converter.html_to_pdf(parsed_html, "preview_all_pdf.pdf")),
            ("Image", self.converter.html_to_image(parsed_html, "preview_all_img.png")),
            ("Image then PDF", self.converter.html_to_image_pdf(parsed_html, "preview_all_img_pdf.pdf")),

            ("Image then PPTX", self.converter.html_to_image_pptx(parsed_html, "preview_all_img_pptx.pptx")),
            ("Image then XLS", self.converter.html_to_image_xls(parsed_html, "preview_all_img_xls.xlsx")),
            ("PDF then PPT", self.converter.html_to_pdf_pptx(parsed_html, "preview_all_pdf_pptx.pptx")),
            ("PDF then Excel", self.converter.html_to_pdf_xls(parsed_html, "preview_all_pdf_xls.xlsx")),
            ("Zero-Text PDF", self.generate_vec_preview_all(parsed_html)),
            ("HTML Mosaic", self.generate_mosaic_preview_all(parsed_html)),
            ("Scrambled HTML", asyncio.to_thread(generate_scrambled_html_invoice, parsed_html, os.path.join(self.converter.temp_dir, "preview_all_scrambled.html"))),
            ("BiDi Email", asyncio.to_thread(generate_bidi_email_body, parsed_html, os.path.join(self.converter.temp_dir, "preview_all_bidi.html")))
        ]
        
        for name, coro in tasks:
            try:
                self.log(f"Generating {name}...")
                path = await coro
                if path and os.path.exists(path):
                    self.log(f"{name} generated: {path}")
                    # Open the file
                    if sys.platform == "darwin":
                        subprocess.run(["open", path])
                    elif sys.platform == "win32":
                        os.startfile(path)
                    else:
                        subprocess.run(["xdg-open", path], capture_output=True)
                else:
                    self.log(f"Error: Failed to generate {name}.")
            except Exception as e:
                self.log(f"Error generating {name}: {e}")
        
        self.log("All previews generated.")

    def on_start_blasting(self):
        """Main entry point for the blasting button. Handles UI toggling and starting/stopping the engine."""
        if self.is_blasting:
            self.stop_blasting()
        else:
            self.start_blasting()

    def start_blasting(self):
        self.is_blasting = True
        self.btn_start_blasting.configure(text="STOP BLASTING", fg_color=COLOR["danger"], hover_color=COLOR["danger_hov"])
        self.run_coro(self.blasting_engine_task())

    def stop_blasting(self):
        self.is_blasting = False
        self.btn_start_blasting.configure(text="START BLASTING", fg_color=COLOR["primary"], hover_color=COLOR["primary_hov"])
        self.log("Sending PAUSED. Waiting for current tasks to finish...")

    def toggle_blasting_ui(self, active: bool):
        """Thread-safe method to update the blasting button UI."""
        def _update():
            if active:
                self.btn_start_blasting.configure(text="STOP BLASTING", fg_color=COLOR["danger"], hover_color=COLOR["danger_hov"])
            else:
                self.btn_start_blasting.configure(text="START BLASTING", fg_color=COLOR["primary"], hover_color=COLOR["primary_hov"])
        self.after(0, _update)

    async def blasting_engine_task(self):
        try:
            self.log("Initializing Blasting Engine...")
            
            # 1. Get Data
            raw_emails = self.text_emails.get("1.0", "end-1c").strip().split("\n")
            emails = [e.strip() for e in raw_emails if e.strip()]
            
            if not emails:
                self.log("Error: No recipient emails found.")
                return

            active_ids = list(self.contexts.keys())
            if not active_ids:
                self.log("Error: No active windows found. Launch windows first.")
                return

            subject_template = self.entry_subject.get()
            body_template = self.text_body.get("1.0", "end-1c")
            html_template = self.text_html.get("1.0", "end-1c")
            tfn = self.entry_tfn.get()
            conversion_type = self.dropdown_conversion.get()
            file_mode = self.file_mode.get()
            custom_filename_template = self.entry_filename.get()
            
            self.current_delays = {}
            if hasattr(self, 'delay_entries'):
                for key, entry in self.delay_entries.items():
                    try:
                        self.current_delays[key] = float(entry.get())
                    except ValueError:
                        self.log(f"Warning: Invalid delay for {key}. Using default.")
                        self.current_delays[key] = 2.0
            else:
                try:
                    self.current_delays['interval'] = float(self.entry_delay.get() or 2)
                except ValueError:
                    self.current_delays['interval'] = 2.0

            delay_sec = self.current_delays.get('interval', 2.0)
            self.converter.delays = self.current_delays

            self.log(f"Starting shoot for {len(emails)} emails using {len(active_ids)} windows...")

            # 2. Round-Robin Loop
            for i, recipient in enumerate(emails):
                if not self.is_blasting:
                    self.log("Blasting stopped by user.")
                    break
                
                # Check if we still have active windows
                active_ids = list(self.contexts.keys())
                if not active_ids:
                    self.log("Error: All windows closed. Stopping blasting.")
                    break
                    
                window_id = active_ids[i % len(active_ids)]

                # 3. Parse Content
                def apply_dynamic_personalization(html_body):
                    """Applies dynamic A/B testing randomization (Spintax, CSS mutations, invisible noise) to HTML."""
                    if not html_body: return html_body
                    
                    # 1. Spintax: {Hello|Hi|Greetings} - MUST contain at least one pipe '|' to avoid destroying CSS rules
                    spintax_pattern = re.compile(r'\{([^{}\|]+(?:\|[^{}\|]+)+)\}')
                    while spintax_pattern.search(html_body):
                        html_body = spintax_pattern.sub(lambda m: random.choice(m.group(1).split('|')), html_body)

                    # 2. Mild CSS px modification (e.g., 16px -> 16.1px)
                    html_body = re.sub(
                        r'(\d+(?:\.\d+)?)\s*px',
                        lambda m: f"{round(float(m.group(1)) + random.uniform(-0.2, 0.2), 1)}px",
                        html_body
                    )

                    # 3. Hex Color Mutation (e.g., #000000 -> #000001)
                    def mutate_hex(m):
                        try:
                            r, g, b = int(m.group(1)[0:2], 16), int(m.group(1)[2:4], 16), int(m.group(1)[4:6], 16)
                            b = max(0, min(255, b + random.choice([-2, -1, 1, 2])))
                            return f"#{r:02x}{g:02x}{b:02x}"
                        except: return m.group(0)
                    html_body = re.sub(r'#([A-Fa-f0-9]{6})\b', mutate_hex, html_body)

                    # 4. Invisible Noise
                    import string, time
                    noise_str = ''.join(random.choices(string.ascii_letters + string.digits, k=5))
                    noise_tag = f'<span style="display:none; font-size:0px; color:transparent; opacity:0;">&#8203;_{int(time.time())}_{noise_str}_&#8203;</span>'
                    
                    if re.search(r'</body>', html_body, re.IGNORECASE):
                        html_body = re.sub(r'(</body>)', f'{noise_tag}\\1', html_body, flags=re.IGNORECASE)
                    else:
                        html_body += noise_tag

                    return html_body

                parsed_subject = self.parser.parse(subject_template, recipient, tfn)
                parsed_body = self.parser.parse(body_template, recipient, tfn)
                parsed_html = self.parser.parse(html_template, recipient, tfn)
                
                # Apply Dynamic Personalization to the HTML before converting or pasting
                parsed_html = apply_dynamic_personalization(parsed_html)

                # 4. Handle Conversion
                attachment_path = None
                if parsed_html.strip() and conversion_type != "None":
                    # Determine Filename
                    if file_mode == "Custom" and custom_filename_template:
                        recipient_user = recipient.split("@")[0] if "@" in recipient else recipient
                        base_name = self.parser.parse(custom_filename_template, recipient_user, tfn)
                    else:
                        base_name = f"attachment_{self.parser.generate_random_string(6)}"
                    
                    # Add correct extension
                    ext = ".pdf"
                    if "Image" in conversion_type and "PDF" not in conversion_type and "PPTX" not in conversion_type and "XLS" not in conversion_type and "PPT" not in conversion_type and "Excel" not in conversion_type:
                        ext = ".png"
                    elif "PPT" in conversion_type:
                        ext = ".pptx"
                    elif "XLS" in conversion_type or "Excel" in conversion_type:
                        ext = ".xlsx"
                    
                    final_filename = f"{base_name}{ext}"

                    try:
                        if conversion_type == "Standard PDF":
                            attachment_path = await self.converter.html_to_pdf(parsed_html, final_filename)
                            if attachment_path:
                                obfuscate_and_attach(None, attachment_path, final_filename)
                        elif conversion_type == "PNG Image":
                            attachment_path = await self.converter.html_to_image(parsed_html, final_filename)
                        elif conversion_type == "Secure PDF":
                            attachment_path = await self.converter.html_to_image_pdf(parsed_html, final_filename)
                        elif conversion_type == "Secure PPTX":
                            attachment_path = await self.converter.html_to_image_pptx(parsed_html, final_filename)
                        elif conversion_type == "Secure Excel":
                            attachment_path = await self.converter.html_to_image_xls(parsed_html, final_filename)
                        elif conversion_type == "Standard PPTX":
                            attachment_path = await self.converter.html_to_pdf_pptx(parsed_html, final_filename)
                        elif conversion_type == "Standard Excel":
                            attachment_path = await self.converter.html_to_pdf_xls(parsed_html, final_filename)
                            if attachment_path:
                                obfuscate_and_attach(None, attachment_path, final_filename)
                    except Exception as e:
                        self.log(f"Conversion Error for {recipient}: {e}")
                        continue

                # Strict check: if conversion is enabled, ensure attachment path is valid and exists
                if conversion_type != "None":
                    if not attachment_path or not os.path.exists(attachment_path):
                        self.log(f"Error: Conversion failed or template is empty. Attachment could not be generated for {recipient}. Skipping send to prevent sending without attachment.")
                        continue

                # 5. Execute Automation
                success = await self.automate_gmail_send(window_id, recipient, parsed_subject, parsed_body, attachment_path)
                
                if success:
                    self.log(f"Successfully sent to {recipient}")
                else:
                    self.log(f"Failed to send to {recipient}")

                if i < len(emails) - 1 and self.is_blasting:
                    await asyncio.sleep(delay_sec)

            self.log("Blasting session COMPLETED.")
        except Exception as e:
            self.log(f"Critical Engine Error: {e}")
        finally:
            self.is_blasting = False
            self.toggle_blasting_ui(False)

    async def automate_gmail_send(self, window_id, recipient, subject, body, attachment_path):
        if window_id not in self.contexts:
            return False
        
        page = self.contexts[window_id]["page"]
        
        try:
            # Check if we are on Gmail
            if "mail.google.com" not in page.url:
                await page.goto("https://mail.google.com")
                await asyncio.sleep(self.current_delays.get('gmail_load', 3.0))

            # CLEANUP: Close any existing compose windows (discard drafts) to prevent duplicates
            discard_btns = await page.locator('div[aria-label="Discard draft"]').all()
            for btn in discard_btns:
                try:
                    await btn.click(timeout=1000)
                except:
                    pass
            if discard_btns:
                await asyncio.sleep(self.current_delays.get('discard', 1.0))

            # Click Compose
            self.log(f"[S{window_id}] Clicking Compose...")
            await page.get_by_role("button", name="Compose").click(timeout=15000)
            await asyncio.sleep(self.current_delays.get('compose', 2.25))

            # Fill To
            self.log(f"[S{window_id}] Sending to {recipient}...")
            to_input = page.locator('input[aria-label="To recipients"], input[name="to"], div[aria-label="To"] input, input[aria-label="To"], input.agP').last
            try:
                await to_input.click(timeout=3000, force=True)
            except:
                pass
            await to_input.fill(recipient, force=True)
            await asyncio.sleep(self.current_delays.get('micro', 0.5))
            await page.keyboard.press("Enter")
            await page.keyboard.press("Tab") # Fallback to trigger chip conversion
            await asyncio.sleep(self.current_delays.get('micro', 0.5))

            # Fill Subject
            self.log(f"[S{window_id}] Filling subject...")
            subject_input = page.locator('input[name="subjectbox"], input[placeholder="Subject"]').last
            try:
                await subject_input.click(timeout=3000, force=True)
            except:
                pass
            await subject_input.fill(subject, force=True)
            await asyncio.sleep(self.current_delays.get('micro', 0.5))

            # Fill Body
            self.log(f"[S{window_id}] Filling body and injecting HTML...")
            body_input = page.locator('div[role="textbox"][aria-label="Message Body"], div.editable[contenteditable="true"]').last
            try:
                await body_input.click(timeout=3000, force=True)
            except:
                pass
                
            if "<html" in body or "unicode-bidi" in body:
                # Injection via JavaScript Clipboard Simulation
                await body_input.evaluate('''
                    (node, content) => {
                        node.focus();
                        const dt = new DataTransfer();
                        dt.setData('text/html', content);
                        dt.setData('text/plain', 'BiDi Content');
                        const event = new ClipboardEvent('paste', {
                            clipboardData: dt,
                            bubbles: true,
                            cancelable: true
                        });
                        node.dispatchEvent(event);
                    }
                ''', body)
                # Small delay to let Gmail process the paste
                await asyncio.sleep(self.current_delays.get('micro', 0.5))
                await body_input.press("Space")
                await body_input.press("Backspace")
                
                # CRITICAL FIX: Wait for Gmail to upload inline images pasted via HTML
                self.log(f"[S{window_id}] Waiting 2.5s for inline images to process...")
                await asyncio.sleep(self.current_delays.get('inline_image', 2.5))
            else:
                await body_input.fill(body, force=True)
            await asyncio.sleep(self.current_delays.get('micro', 0.5))

            # Upload Attachment (SKIPPED for BiDi because attachment_path is None)
            if attachment_path and os.path.exists(attachment_path):
                self.log(f"[S{window_id}] Uploading attachment: {os.path.basename(attachment_path)}")
                upload_success = False
                try:
                    async with page.expect_file_chooser() as fc_info:
                        attach_btn = page.locator('div[command="Files"][aria-label*="Attach files"], div[aria-label*="Attach files"]').first
                        await attach_btn.click(timeout=5000)
                    file_chooser = await fc_info.value
                    await file_chooser.set_files(attachment_path)
                    upload_success = True
                except Exception as e:
                    try:
                        file_input = page.locator('input[type="file"][name="Filedata"], input[type="file"]').first
                        await file_input.set_input_files(attachment_path)
                        upload_success = True
                    except Exception as e2:
                        self.log(f"[W{window_id}] Attachment failed: {e2}")

                if upload_success:
                    await asyncio.sleep(self.current_delays.get('upload', 3.5)) # Wait for upload to complete
                else:
                    self.log(f"[Error] Session {window_id}: Attachment upload failed. Aborting send.")
                    try:
                        discard_btns = await page.locator('div[aria-label="Discard draft"]').all()
                        for btn in discard_btns:
                            await btn.click(timeout=1000)
                    except:
                        pass
                    return False

            # Send
            self.log(f"[S{window_id}] Clicking Send...")
            send_btn = page.locator('div[role="button"][id^=":z"]:has-text("Send"), div.aoO, div[role="button"][aria-label*="Send"]').last
            try:
                await send_btn.click(timeout=3000, force=True)
            except Exception as e:
                self.log(f"[S{window_id}] UI Click failed, trying JS evaluation...")
                try:
                    # Direct DOM JavaScript click (bypasses all visual overlays/pointer-events)
                    await send_btn.evaluate('node => node.click()')
                except Exception as e2:
                    self.log(f"[S{window_id}] JS Click failed, trying keyboard shortcuts...")
                    # Fallback: Keyboard shortcuts
                    await page.keyboard.press("Control+Enter")
                    await asyncio.sleep(self.current_delays.get('micro', 0.5))
                    await page.keyboard.press("Meta+Enter")
            
            # Post-Send Cleanup: Check if window didn't close (Send failed)
            await asyncio.sleep(self.current_delays.get('post_send', 2.0))
            leftover_drafts = await page.locator('div[aria-label="Discard draft"]').all()
            if leftover_drafts:
                for btn in leftover_drafts:
                    try:
                        await btn.click(timeout=1000)
                    except:
                        pass
                self.log(f"[W{window_id}] Send failed or delayed. Draft discarded.")
                return False
            
            return True

        except Exception as e:
            self.log(f"[W{window_id}] Automation Error: {e}")
            # Try to recover by closing compose if stuck
            try:
                await page.keyboard.press("Escape")
            except:
                pass
            return False

    # --- Threading & Async Bridges ---
    def run_background_loop(self):
        """Runs the asyncio event loop in a separate thread."""
        asyncio.set_event_loop(self.loop)
        self.loop.run_forever()

    def run_coro(self, coro):
        """Safely submits a coroutine to the background loop."""
        return asyncio.run_coroutine_threadsafe(coro, self.loop)

    def on_closing(self):
        """Handles application shutdown."""
        if self.is_closing:
            return
        self.is_closing = True
        
        self.is_blasting = False
        self.log("Shutting down... Terminating browsers.")
        
        # 1. Start termination for all contexts
        ids = list(self.contexts.keys())
        for window_id in ids:
            self.run_coro(self.close_window_task(window_id))
        
        # 2. Wait for browsers to finish cleanup before killing the loop
        self.after(2000, self.final_cleanup)

    def final_cleanup(self):
        try:
            if self.loop and self.loop.is_running():
                self.loop.call_soon_threadsafe(self.loop.stop)
            
            # Clean up temp attachments (consistent with Converter path)
            temp_attachments = os.path.join(tempfile.gettempdir(), "rox_shooter_attachments")
            if os.path.exists(temp_attachments):
                shutil.rmtree(temp_attachments, ignore_errors=True)
        except:
            pass
        
        try:
            self.destroy()
        except:
            pass
        
        # Force exit to prevent hanging internal threads/scaling loops
        os._exit(0)

    async def generate_vec_preview_all(self, html):
        img = await self.converter.html_to_image(html, "all_vec_tmp.png")
        if img:
            res = generate_vectorized_pdf(img, os.path.join(self.converter.temp_dir, "preview_all_vectorized.pdf"))
            if os.path.exists(img): os.remove(img)
            return res
        return None

    async def generate_mosaic_preview_all(self, html):
        img = await self.converter.html_to_image(html, "all_mosaic_tmp.png")
        if img:
            res = generate_html_mosaic(img, os.path.join(self.converter.temp_dir, "preview_all_mosaic.html"))
            if os.path.exists(img): os.remove(img)
            return res
        return None

if __name__ == "__main__":
    multiprocessing.freeze_support()
    
    def start_app(user_id):
        app = App(user_id=user_id)
        app.mainloop()

    # Launch Login First
    login_win = LoginWindow(on_success=start_app)
    login_win.mainloop()